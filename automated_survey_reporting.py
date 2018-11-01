#By the help of xlsxwriter, random and openpyxl this code
    #creates an Excel file
    #fills out columns with the answers of an imaginery satisfaction survey
    #summarizes and visualizes the results
    #creates a Power Point Presentation with a straightforward summing-up of the results.



import xlsxwriter
import random
import openpyxl
from openpyxl import load_workbook
import pptx
from pptx import Presentation



workbook = xlsxwriter.Workbook('Survey_assessment.xlsx')
worksheet = workbook.add_worksheet()

worksheet.write('A1','Gender')
worksheet.write('B1','Age')
worksheet.write('C1','Satisfaction')

genders = workbook.add_chart({'type': 'pie'})

satisf = workbook.add_chart({'type': 'column'})

agegroups = workbook.add_chart({'type': 'column'})


b=['male','female']
for row in range(1,101):
    worksheet.write(row,0,random.choice(b))
for row in range(1,101):
    worksheet.write(row,1,random.randint(18,67))
for row in range(1,101):
    worksheet.write(row,2,random.randint(1,5))


worksheet.write('E1', 'Genders')
worksheet.write('E2', 'Females')
worksheet.write('E3', 'Males')
worksheet.write('H1', 'Age groups')
worksheet.write('H2', '18-27')
worksheet.write('H3', '28-37')
worksheet.write('H4', '38-47')
worksheet.write('H5', '48-57')
worksheet.write('H6', '58-67')
worksheet.write('K1', 'Satisfaction')
worksheet.write('K2', 'Unsatisfied')
worksheet.write('K3', 'Somewhat unsatisfied')
worksheet.write('K4', 'Indifferent')
worksheet.write('K5', 'Somewhat satisfied')
worksheet.write('K6', 'Satisfied')
worksheet.write('N1', 'Average age')




genders.add_series({
    'name': 'Genders',
    'categories': '=Sheet1!$E$2:$E$3',
    'values':     '=Sheet1!$F$2:$F$3',
    'data_labels': {'value': True},
})

satisf.add_series({
    'name': 'Age groups',
    'categories': '=Sheet1!$H$2:$H$6',
    'values':     '=Sheet1!$I$2:$I$6',
    'data_labels': {'value': True},
    'fill':   {'color': 'green'},
})

agegroups.add_series({
    'name': 'Satisfaction',
    'categories': '=Sheet1!$K$2:$K$6',
    'values':     '=Sheet1!$L$2:$L$6',
    'data_labels': {'value': True},
    'fill':   {'color': 'orange'},
})



worksheet.insert_chart('D10', genders)
worksheet.insert_chart('L10', satisf)
worksheet.insert_chart('T10', agegroups)


workbook.close()

wb=load_workbook('Survey_assessment.xlsx')
sheet=wb.active

female=0
male=0
agegroup1=0
agegroup2=0
agegroup3=0
agegroup4=0
agegroup5=0
unsatisfied=0
swunsatisfied=0
indifferent=0
swsatisfied=0
satisfied=0
avgage=0


for row in sheet.iter_rows(min_row=2, max_row=101, min_col=1, max_col=1):
    for cell in row:
        if cell.value=="female":
            female+=1
        else:
            male+=1

for row in sheet.iter_rows(min_row=2, max_row=101, min_col=2, max_col=2):
    for cell in row:
        if 18<=cell.value<=27:
            agegroup1+=1
        elif 28<=cell.value<=37:
            agegroup2+=1
        elif 38<=cell.value<=47:
            agegroup3+=1
        elif 48<=cell.value<=57:
            agegroup4+=1
        else:
            agegroup5+=1

for row in sheet.iter_rows(min_row=2, max_row=101, min_col=3, max_col=3):
    for cell in row:
        if cell.value==1:
            unsatisfied+=1
        elif cell.value==2:
            swunsatisfied+=1
        elif cell.value==3:
            indifferent+=1
        elif cell.value==4:
            swsatisfied+=1
        else:
            satisfied+=1

for row in sheet.iter_rows(min_row=2, max_row=101, min_col=2, max_col=2):
    for cell in row:
        avgage+=cell.value

sheet.cell(row=2, column=6).value=female
sheet.cell(row=3, column=6).value=male

sheet.cell(row=2, column=9).value=agegroup1
sheet.cell(row=3, column=9).value=agegroup2
sheet.cell(row=4, column=9).value=agegroup3
sheet.cell(row=5, column=9).value=agegroup4
sheet.cell(row=6, column=9).value=agegroup5

sheet.cell(row=2, column=12).value=unsatisfied
sheet.cell(row=3, column=12).value=swunsatisfied
sheet.cell(row=4, column=12).value=indifferent
sheet.cell(row=5, column=12).value=swsatisfied
sheet.cell(row=6, column=12).value=satisfied
sheet.cell(row=2, column=14).value=(avgage/100)

gender1=sheet['F2'].value
gender2=sheet['F3'].value
averageage=sheet['N2'].value

wb.save('Survey_assessment.xlsx')
wb.close()

prezi=Presentation()
title_slide_layout = prezi.slide_layouts[0]
slide = prezi.slides.add_slide(title_slide_layout)

title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Survey assessment"
subtitle.text = "Created by Andras Komjathy"

overviewslide = prezi.slide_layouts[1]

slide = prezi.slides.add_slide(overviewslide)
shapes = slide.shapes

title_shape = shapes.title
body_shape = shapes.placeholders[1]

title_shape.text = "Overview"
tf=body_shape.text_frame
tf.text='Number of participants=100'
tf1=tf.add_paragraph()
tf1.text='Genders'
p=tf.add_paragraph()
p.text='Females: {}'.format(gender1)
p.level=1

p2=tf.add_paragraph()
p2.text='Males: {}'.format(gender2)
p2.level=1

tf2=tf.add_paragraph()
tf2.text='Average age: {}'.format(averageage)


agegroupslide = prezi.slide_layouts[1]
slide = prezi.slides.add_slide(agegroupslide)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Participants in age groups"
tf=body_shape.text_frame
tf.text='18-27: {}'.format(agegroup1)
tf1=tf.add_paragraph()
tf1.text='28-37: {}'.format(agegroup2)
tf2=tf.add_paragraph()
tf2.text='38-47: {}'.format(agegroup3)
tf3=tf.add_paragraph()
tf3.text='48-57: {}'.format(agegroup4)
tf4=tf.add_paragraph()
tf4.text='58-67: {}'.format(agegroup5)


satisfactionslide = prezi.slide_layouts[1]
slide = prezi.slides.add_slide(satisfactionslide)
shapes = slide.shapes
title_shape = shapes.title
body_shape = shapes.placeholders[1]
title_shape.text = "Participants' satisfaction"
tf=body_shape.text_frame
tf.text='Unsatisfied: {}'.format(unsatisfied)
tf1=tf.add_paragraph()
tf1.text='Somewhat unsatisfied: {}'.format(swunsatisfied)
tf2=tf.add_paragraph()
tf2.text='Indifferent: {}'.format(indifferent)
tf3=tf.add_paragraph()
tf3.text='Somewhat satisfied: {}'.format(swsatisfied)
tf4=tf.add_paragraph()
tf4.text='Satisfied: {}'.format(satisfied)



prezi.save('Survey_assessment.pptx')






